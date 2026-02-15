"""
The Librarian — Document Format Readers (Phase 12)

Read-on-demand handlers for rich document formats.
All dependencies are optional — graceful degradation with install hints.

Supported formats:
    - PDF  (.pdf)   via PyMuPDF (fitz)
    - DOCX (.docx)  via python-docx
    - XLSX (.xlsx)   via openpyxl
    - PPTX (.pptx)  via python-pptx
    - TXT  (.txt, .md, .csv, .json, etc.) via built-in open()
"""
import hashlib
import os
from dataclasses import dataclass, field
from typing import Optional, List, Dict, Any


# ─── Result Types ─────────────────────────────────────────────────────────────

@dataclass
class DocumentMetadata:
    """Metadata extracted from a document during registration."""
    file_name: str = ""
    file_path: str = ""
    file_type: str = ""
    file_hash: str = ""
    title: str = ""
    page_count: Optional[int] = None
    author: Optional[str] = None
    file_size: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ReadResult:
    """Result of reading a document."""
    success: bool = False
    text: str = ""
    pages: Optional[List[str]] = None  # Per-page text (if applicable)
    headings: Optional[List[str]] = None  # Extracted headings/structure
    error: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)


# ─── File Hash ────────────────────────────────────────────────────────────────

def compute_file_hash(file_path: str) -> str:
    """Compute SHA-256 hash of a file for change detection."""
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            sha256.update(chunk)
    return sha256.hexdigest()


# ─── Format Detection ────────────────────────────────────────────────────────

SUPPORTED_TYPES = {"pdf", "docx", "xlsx", "pptx", "txt"}

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".conf", ".log", ".rst", ".html",
    ".htm", ".py", ".js", ".ts", ".sh", ".bash", ".sql",
}


def detect_file_type(file_path: str) -> str:
    """Detect document type from file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    type_map = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".xlsx": "xlsx",
        ".pptx": "pptx",
    }
    if ext in type_map:
        return type_map[ext]
    if ext in TEXT_EXTENSIONS:
        return "txt"
    return "unknown"


# ─── PDF Reader ──────────────────────────────────────────────────────────────

def _read_pdf(file_path: str, pages: Optional[str] = None) -> ReadResult:
    """Read a PDF file using PyMuPDF (fitz)."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return ReadResult(
            success=False,
            error="PyMuPDF not installed. Install with: pip install PyMuPDF --break-system-packages"
        )

    try:
        doc = fitz.open(file_path)
        total_pages = len(doc)

        # Parse page range
        page_indices = _parse_page_range(pages, total_pages) if pages else range(total_pages)

        page_texts = []
        for i in page_indices:
            if 0 <= i < total_pages:
                page = doc[i]
                text = page.get_text("text")
                page_texts.append(text.strip())

        doc.close()

        full_text = "\n\n".join(
            f"[Page {i+1}]\n{text}" for i, text in zip(page_indices, page_texts) if text
        )

        return ReadResult(
            success=True,
            text=full_text,
            pages=page_texts,
            metadata={"total_pages": total_pages, "pages_read": len(page_texts)},
        )
    except Exception as e:
        return ReadResult(success=False, error=f"PDF read error: {e}")


def _get_pdf_metadata(file_path: str) -> DocumentMetadata:
    """Extract PDF metadata without reading full content."""
    meta = DocumentMetadata(
        file_name=os.path.basename(file_path),
        file_path=file_path,
        file_type="pdf",
        file_size=os.path.getsize(file_path),
    )
    try:
        import fitz
        doc = fitz.open(file_path)
        meta.page_count = len(doc)
        pdf_meta = doc.metadata
        if pdf_meta:
            meta.title = pdf_meta.get("title", "") or ""
            meta.author = pdf_meta.get("author", "") or ""
            meta.metadata = {k: v for k, v in pdf_meta.items() if v}
        doc.close()
    except ImportError:
        meta.metadata["note"] = "PyMuPDF not installed — metadata extraction limited"
    except Exception as e:
        meta.metadata["error"] = str(e)
    return meta


# ─── DOCX Reader ─────────────────────────────────────────────────────────────

def _read_docx(file_path: str, pages: Optional[str] = None) -> ReadResult:
    """Read a DOCX file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        return ReadResult(
            success=False,
            error="python-docx not installed. Install with: pip install python-docx --break-system-packages"
        )

    try:
        doc = Document(file_path)
        paragraphs = []
        headings = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading", "").strip() or "1"
                headings.append(f"H{level}: {text}")
                paragraphs.append(f"\n[{para.style.name}] {text}")
            else:
                paragraphs.append(text)

        # Also extract table content
        for table in doc.tables:
            table_rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_rows.append(" | ".join(cells))
            if table_rows:
                paragraphs.append("\n[Table]\n" + "\n".join(table_rows))

        full_text = "\n".join(paragraphs)

        return ReadResult(
            success=True,
            text=full_text,
            headings=headings,
            metadata={
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "heading_count": len(headings),
            },
        )
    except Exception as e:
        return ReadResult(success=False, error=f"DOCX read error: {e}")


def _get_docx_metadata(file_path: str) -> DocumentMetadata:
    """Extract DOCX metadata."""
    meta = DocumentMetadata(
        file_name=os.path.basename(file_path),
        file_path=file_path,
        file_type="docx",
        file_size=os.path.getsize(file_path),
    )
    try:
        from docx import Document
        doc = Document(file_path)
        props = doc.core_properties
        meta.title = props.title or ""
        meta.author = props.author or ""
        meta.page_count = len(doc.paragraphs)  # Paragraph count as proxy
        meta.metadata = {
            "paragraph_count": len(doc.paragraphs),
            "table_count": len(doc.tables),
        }
        if props.subject:
            meta.metadata["subject"] = props.subject
    except ImportError:
        meta.metadata["note"] = "python-docx not installed — metadata extraction limited"
    except Exception as e:
        meta.metadata["error"] = str(e)
    return meta


# ─── XLSX Reader ─────────────────────────────────────────────────────────────

def _read_xlsx(file_path: str, pages: Optional[str] = None) -> ReadResult:
    """Read an XLSX file using openpyxl. 'pages' is interpreted as sheet names/indices."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ReadResult(
            success=False,
            error="openpyxl not installed. Install with: pip install openpyxl --break-system-packages"
        )

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        sheet_names = wb.sheetnames

        # Filter sheets if specified
        target_sheets = sheet_names
        if pages:
            target_sheets = [s for s in pages.split(",") if s.strip() in sheet_names]
            if not target_sheets:
                target_sheets = sheet_names  # Fallback to all

        sheet_texts = []
        for sheet_name in target_sheets:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                sheet_texts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))

        wb.close()
        full_text = "\n\n".join(sheet_texts)

        return ReadResult(
            success=True,
            text=full_text,
            metadata={"sheet_count": len(sheet_names), "sheets": sheet_names},
        )
    except Exception as e:
        return ReadResult(success=False, error=f"XLSX read error: {e}")


def _get_xlsx_metadata(file_path: str) -> DocumentMetadata:
    """Extract XLSX metadata."""
    meta = DocumentMetadata(
        file_name=os.path.basename(file_path),
        file_path=file_path,
        file_type="xlsx",
        file_size=os.path.getsize(file_path),
    )
    try:
        from openpyxl import load_workbook
        wb = load_workbook(file_path, read_only=True)
        meta.page_count = len(wb.sheetnames)
        meta.metadata = {"sheets": wb.sheetnames}
        wb.close()
    except ImportError:
        meta.metadata["note"] = "openpyxl not installed — metadata extraction limited"
    except Exception as e:
        meta.metadata["error"] = str(e)
    return meta


# ─── PPTX Reader ─────────────────────────────────────────────────────────────

def _read_pptx(file_path: str, pages: Optional[str] = None) -> ReadResult:
    """Read a PPTX file using python-pptx. 'pages' is interpreted as slide numbers."""
    try:
        from pptx import Presentation
    except ImportError:
        return ReadResult(
            success=False,
            error="python-pptx not installed. Install with: pip install python-pptx --break-system-packages"
        )

    try:
        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        # Parse slide range
        slide_indices = _parse_page_range(pages, total_slides) if pages else range(total_slides)

        slide_texts = []
        for i in slide_indices:
            if 0 <= i < total_slides:
                slide = prs.slides[i]
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                texts.append(text)
                if texts:
                    slide_texts.append(f"[Slide {i+1}]\n" + "\n".join(texts))

        full_text = "\n\n".join(slide_texts)

        return ReadResult(
            success=True,
            text=full_text,
            pages=[t for t in slide_texts],
            metadata={"total_slides": total_slides, "slides_read": len(slide_texts)},
        )
    except Exception as e:
        return ReadResult(success=False, error=f"PPTX read error: {e}")


def _get_pptx_metadata(file_path: str) -> DocumentMetadata:
    """Extract PPTX metadata."""
    meta = DocumentMetadata(
        file_name=os.path.basename(file_path),
        file_path=file_path,
        file_type="pptx",
        file_size=os.path.getsize(file_path),
    )
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        meta.page_count = len(prs.slides)
        meta.metadata = {"slide_count": len(prs.slides)}
    except ImportError:
        meta.metadata["note"] = "python-pptx not installed — metadata extraction limited"
    except Exception as e:
        meta.metadata["error"] = str(e)
    return meta


# ─── Text Reader ─────────────────────────────────────────────────────────────

def _read_text(file_path: str, pages: Optional[str] = None) -> ReadResult:
    """Read a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return ReadResult(
            success=True,
            text=text,
            metadata={"char_count": len(text), "line_count": text.count("\n") + 1},
        )
    except Exception as e:
        return ReadResult(success=False, error=f"Text read error: {e}")


def _get_text_metadata(file_path: str) -> DocumentMetadata:
    """Extract text file metadata."""
    return DocumentMetadata(
        file_name=os.path.basename(file_path),
        file_path=file_path,
        file_type="txt",
        file_size=os.path.getsize(file_path),
    )


# ─── Unified Interface ───────────────────────────────────────────────────────

def read_document(file_path: str, pages: Optional[str] = None) -> ReadResult:
    """Read a document using the appropriate format handler.

    Args:
        file_path: Absolute path to the document.
        pages: Optional page/sheet range (e.g., "1-5", "Sheet1,Sheet2").

    Returns:
        ReadResult with extracted text and metadata.
    """
    if not os.path.isfile(file_path):
        return ReadResult(success=False, error=f"File not found: {file_path}")

    file_type = detect_file_type(file_path)
    readers = {
        "pdf": _read_pdf,
        "docx": _read_docx,
        "xlsx": _read_xlsx,
        "pptx": _read_pptx,
        "txt": _read_text,
    }

    reader = readers.get(file_type)
    if not reader:
        return ReadResult(
            success=False,
            error=f"Unsupported file type: {file_type} (extension: {os.path.splitext(file_path)[1]})"
        )

    return reader(file_path, pages)


def get_document_metadata(file_path: str) -> DocumentMetadata:
    """Extract metadata from a document without reading full content.

    Args:
        file_path: Absolute path to the document.

    Returns:
        DocumentMetadata with file info, title, page count, etc.
    """
    if not os.path.isfile(file_path):
        return DocumentMetadata(metadata={"error": f"File not found: {file_path}"})

    file_type = detect_file_type(file_path)
    extractors = {
        "pdf": _get_pdf_metadata,
        "docx": _get_docx_metadata,
        "xlsx": _get_xlsx_metadata,
        "pptx": _get_pptx_metadata,
        "txt": _get_text_metadata,
    }

    extractor = extractors.get(file_type)
    if not extractor:
        return DocumentMetadata(
            file_name=os.path.basename(file_path),
            file_path=file_path,
            file_type=file_type,
            file_size=os.path.getsize(file_path),
            metadata={"note": f"No metadata extractor for type: {file_type}"},
        )

    meta = extractor(file_path)
    meta.file_hash = compute_file_hash(file_path)
    return meta


# ─── Helpers ──────────────────────────────────────────────────────────────────

def _parse_page_range(pages_str: str, total: int) -> List[int]:
    """Parse a page range string like '1-5' or '1,3,5' into 0-based indices."""
    indices = []
    for part in pages_str.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            try:
                s = max(1, int(start))
                e = min(total, int(end))
                indices.extend(range(s - 1, e))  # Convert to 0-based
            except ValueError:
                continue
        else:
            try:
                idx = int(part)
                if 1 <= idx <= total:
                    indices.append(idx - 1)  # Convert to 0-based
            except ValueError:
                continue
    return indices if indices else list(range(total))
